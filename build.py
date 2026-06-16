#!/usr/bin/env python3
"""
build.py - keep the HTML viewer in sync with the PowerPoint file.

What it does:
  1. Converts presentation/Welcome_Inbet_Online.pptx -> PDF   (LibreOffice)
  2. Renders each page to slides/slide-NN.jpg                 (Poppler: pdftoppm)
  3. Optimises the images                                     (Pillow)
  4. Extracts the Inbet logo and every slide's hyperlinks from the .pptx
  5. Writes manifest.json + manifest.js (the slide list the viewer reads)

Run it whenever you change the .pptx:
    python build.py

All tools are free / open-source:
  - LibreOffice  (soffice)      https://www.libreoffice.org/
  - Poppler      (pdftoppm)     https://poppler.freedesktop.org/
  - Pillow       (pip install Pillow)
Optional, for nicer auto-titles:
  - python-pptx  (pip install python-pptx)

Note on fonts: the deck's heading font is "Trebuchet MS". On Linux (incl. the
GitHub Actions runner) that font is usually absent, so LibreOffice substitutes a
wider one and long titles wrap and get clipped. To prevent that we alias
Trebuchet MS -> Carlito (a free, narrower Calibri-metric font) just for the
render. Install it with:  sudo apt-get install fonts-crosextra-carlito
On Windows/macOS, where Trebuchet MS is present, this step is skipped.
"""

import glob
import html
import io
import json
import os
import re
import shutil
import subprocess
import sys
import zipfile

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, "presentation", "Welcome_Inbet_Online.pptx")
SLIDES_DIR = os.path.join(HERE, "slides")
DPI = 150
JPEG_QUALITY = 88


def find_tool(*names):
    for n in names:
        p = shutil.which(n)
        if p:
            return p
    return None


def run(cmd, env=None):
    print("  $", " ".join(cmd))
    subprocess.run(cmd, check=True, env=env)


def normalize_dashes(s):
    """Replace em/en dashes (and surrounding spaces) with a plain hyphen."""
    return re.sub(r"\s*[\u2012\u2013\u2014\u2015]\s*", " - ", s).strip()


def render_env():
    """On Linux, return an env that makes LibreOffice substitute the missing
    Trebuchet MS with Carlito so titles render on one line. No-op elsewhere."""
    if not sys.platform.startswith("linux"):
        return None
    conf = os.path.join(HERE, ".fonts-render.conf")
    with open(conf, "w", encoding="utf-8") as f:
        f.write(
            '<?xml version="1.0"?>\n'
            '<!DOCTYPE fontconfig SYSTEM "fonts.dtd">\n'
            "<fontconfig>\n"
            '  <include ignore_missing="yes">/etc/fonts/fonts.conf</include>\n'
            "  <match target=\"pattern\">\n"
            '    <test name="family"><string>Trebuchet MS</string></test>\n'
            '    <edit name="family" mode="assign" binding="strong">'
            "<string>Carlito</string></edit>\n"
            "  </match>\n"
            "</fontconfig>\n"
        )
    env = dict(os.environ)
    env["FONTCONFIG_FILE"] = conf
    return env


def existing_titles():
    """Preserve hand-edited titles from a previous manifest.json when possible."""
    path = os.path.join(HERE, "manifest.json")
    if not os.path.exists(path):
        return []
    try:
        data = json.load(open(path, encoding="utf-8"))
        return [s.get("title", "") for s in data.get("slides", [])]
    except Exception:
        return []


def titles_from_pptx(n_slides):
    """Best-effort: first text line of each slide. Falls back to 'Slide N'."""
    try:
        from pptx import Presentation
    except Exception:
        return ["Slide %d" % (i + 1) for i in range(n_slides)]
    titles = []
    prs = Presentation(PPTX)
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    line = txt.splitlines()[0].strip()
                    if 2 <= len(line) <= 42:
                        title = line
                        break
        titles.append(title or ("Slide %d" % (i + 1)))
    return titles


def links_per_slide(n_slides):
    """Pull external hyperlinks (label + url) out of every slide, in order."""
    result = [[] for _ in range(n_slides)]
    try:
        with zipfile.ZipFile(PPTX) as z:
            for sn in range(1, n_slides + 1):
                sx = "ppt/slides/slide%d.xml" % sn
                rx = "ppt/slides/_rels/slide%d.xml.rels" % sn
                if sx not in z.namelist():
                    continue
                xml = z.read(sx).decode("utf-8", "ignore")
                rels = z.read(rx).decode("utf-8", "ignore") if rx in z.namelist() else ""
                rid = {}
                for m in re.finditer(r"<Relationship\b[^>]*>", rels):
                    i = re.search(r'Id="([^"]+)"', m.group(0))
                    t = re.search(r'Target="([^"]+)"', m.group(0))
                    if i and t:
                        rid[i.group(1)] = html.unescape(t.group(1))
                out, seen = [], set()
                for para in re.finditer(r"<a:p>.*?</a:p>", xml, re.S):
                    cur, buf = None, ""
                    for run_m in re.finditer(r"<a:r>.*?</a:r>", para.group(0), re.S):
                        r = run_m.group(0)
                        hl = re.search(r'<a:hlinkClick[^>]*r:id="([^"]+)"', r)
                        t = re.search(r"<a:t>(.*?)</a:t>", r, re.S)
                        txt = html.unescape(t.group(1)) if t else ""
                        if hl:
                            if cur == hl.group(1):
                                buf += txt
                            else:
                                if cur:
                                    out.append((buf.strip(), rid.get(cur, "")))
                                cur, buf = hl.group(1), txt
                        else:
                            if cur:
                                out.append((buf.strip(), rid.get(cur, "")))
                            cur, buf = None, ""
                    if cur:
                        out.append((buf.strip(), rid.get(cur, "")))
                for label, url in out:
                    if url.startswith("http") and label and url not in seen:
                        seen.add(url)
                        result[sn - 1].append({"label": label, "url": url})
    except Exception as e:
        print("     link extraction skipped (%s)" % e)
    return result


def extract_logo():
    """Best-effort: pull the wordmark logo out of the .pptx media into
    assets/inbet-logo.png. Never fatal - keeps the existing asset on failure."""
    try:
        from PIL import Image
        os.makedirs(os.path.join(HERE, "assets"), exist_ok=True)
        best = None  # (area, name, bytes)
        with zipfile.ZipFile(PPTX) as z:
            for n in z.namelist():
                if not n.lower().startswith("ppt/media/") or not n.lower().endswith(".png"):
                    continue
                data = z.read(n)
                im = Image.open(io.BytesIO(data))
                w, h = im.size
                has_alpha = im.mode in ("RGBA", "LA") or "transparency" in im.info
                if has_alpha and w > h * 1.8 and (w * h) < 250000:  # wide wordmark
                    area = w * h
                    if best is None or area < best[0]:
                        best = (area, n, data)
        if best:
            with open(os.path.join(HERE, "assets", "inbet-logo.png"), "wb") as f:
                f.write(best[2])
            print("     logo extracted from %s" % best[1])
        else:
            print("     no logo candidate found - keeping existing asset")
    except Exception as e:
        print("     logo extraction skipped (%s) - keeping existing asset" % e)


def main():
    if not os.path.exists(PPTX):
        sys.exit("ERROR: %s not found." % PPTX)

    soffice = find_tool("soffice", "libreoffice")
    pdftoppm = find_tool("pdftoppm")
    if not soffice:
        sys.exit("ERROR: LibreOffice (soffice) not found on PATH.")
    if not pdftoppm:
        sys.exit("ERROR: Poppler (pdftoppm) not found on PATH.")
    try:
        from PIL import Image
    except Exception:
        sys.exit("ERROR: Pillow not installed.  Run: pip install Pillow")

    env = render_env()
    if env and not (find_tool("fc-list")
                    and "carlito" in subprocess.run(["fc-list"], capture_output=True,
                                                     text=True).stdout.lower()):
        print("     NOTE: Carlito font not found; titles may wrap. "
              "Install with: sudo apt-get install fonts-crosextra-carlito")

    print("1/4  Converting PPTX -> PDF ...")
    run([soffice, "--headless", "--convert-to", "pdf", "--outdir", HERE, PPTX], env=env)
    pdf = os.path.join(HERE, "Welcome_Inbet_Online.pdf")

    print("2/4  Rendering PDF pages -> PNG ...")
    os.makedirs(SLIDES_DIR, exist_ok=True)
    for f in glob.glob(os.path.join(SLIDES_DIR, "slide-*")):
        os.remove(f)
    run([pdftoppm, "-png", "-r", str(DPI), pdf, os.path.join(SLIDES_DIR, "slide")])

    print("3/4  Optimising images -> JPEG ...")
    for p in sorted(glob.glob(os.path.join(SLIDES_DIR, "slide-*.png"))):
        Image.open(p).convert("RGB").save(
            p[:-4] + ".jpg", "JPEG", quality=JPEG_QUALITY, optimize=True, progressive=True
        )
        os.remove(p)
    os.remove(pdf)
    jpgs = sorted(glob.glob(os.path.join(SLIDES_DIR, "slide-*.jpg")))
    n = len(jpgs)
    print("     %d slides rendered." % n)

    print("     Extracting logo ...")
    extract_logo()

    print("4/4  Writing manifest ...")
    prev = existing_titles()
    if len(prev) == n and all(prev):
        titles = prev
        print("     reused titles from existing manifest.json")
    else:
        titles = titles_from_pptx(n)
        print("     generated titles from the .pptx")
    titles = [normalize_dashes(t) for t in titles]   # em/en dash -> hyphen

    links = links_per_slide(n)
    nlinks = sum(len(x) for x in links)
    print("     %d hyperlinks captured" % nlinks)

    slides = []
    for i, j in enumerate(jpgs):
        entry = {"file": "slides/%s" % os.path.basename(j), "title": titles[i]}
        if links[i]:
            entry["links"] = links[i]
        slides.append(entry)

    manifest = {
        "source": "presentation/Welcome_Inbet_Online.pptx",
        "generated": "build.py output - edit titles freely; re-run build.py to refresh",
        "slides": slides,
    }
    with open(os.path.join(HERE, "manifest.json"), "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)
    with open(os.path.join(HERE, "manifest.js"), "w", encoding="utf-8") as f:
        f.write("// Auto-generated by build.py - slide list the viewer reads.\n")
        f.write("// Loaded via <script src> so index.html works even via file://.\n")
        f.write("window.INBET_SLIDES = " + json.dumps(slides, indent=2, ensure_ascii=False) + ";\n")

    print("\nDone. Open index.html (via Live Server / http) to view %d slides." % n)


if __name__ == "__main__":
    main()
